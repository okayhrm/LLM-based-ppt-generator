import openai
import re
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from duckduckgo_search import DDGS
import tempfile
import time
import json
from typing import List, Tuple, Dict
import os
from PIL import Image
from dotenv import load_dotenv

MODELS = {
    "Mistral 7B": "mistralai/mistral-7b-instruct",
    "Mixtral 8x7B": "mistralai/mixtral-8x7b-instruct",
    "GPT-3.5 Turbo": "openai/gpt-3.5-turbo",
    "Claude 2.1": "anthropic/claude-2.1"
}
DEFAULT_SLIDES = 7
MAX_SLIDES = 20

load_dotenv()

openai.api_key = os.getenv("OPENROUTER_API_KEY")
openai.api_base = "https://openrouter.ai/api/v1"

TEMPLATE_DIR = "templates"
THUMBNAIL_DIR = os.path.join(TEMPLATE_DIR, "thumbnails")

template_files = [f for f in os.listdir(TEMPLATE_DIR) if f.endswith(".pptx")]
template_labels = [f.replace("_", " ").replace(".pptx", "").title() for f in template_files]
template_map = dict(zip(template_labels, template_files))

def extract_slide_count(prompt: str) -> int:
    match = re.search(r'(\d+)[- ]*slide', prompt.lower())
    if match:
        count = int(match.group(1))
        return min(max(1, count), MAX_SLIDES)
    return DEFAULT_SLIDES

def fetch_search_snippets(query: str, count: int = 5) -> List[str]:
    try:
        with DDGS() as ddgs:
            results = ddgs.text(query, max_results=count)
            return [r["body"] for r in results if "body" in r][:count]
    except Exception:
        return []

def generate_slide_content(prompt: str, model: str, use_live_info: bool) -> List[Tuple[str, List[str]]]:
    slide_count = extract_slide_count(prompt)
    clean_topic = re.sub(r'\d+[- ]*slide[s]*', '', prompt, flags=re.IGNORECASE).strip()

    system_prompt = (
        "You are an expert presentation creator. Generate exactly {count} slides in JSON format: "
        "{{\"slides\": [{{\"title\": \"Slide 1 Title\", \"content\": [\"Point 1\", \"Point 2\", \"Point 3\"]}}]}}\n"
        "Rules:\n"
        "1. Each slide has 1 title and 3-5 bullet points\n"
        "2. Use professional business language\n"
        "3. Content must be factual and educational\n"
        "4. Final output MUST be valid JSON only"
    ).format(count=slide_count)

    user_prompt = f"Create presentation about: {clean_topic}"
    if use_live_info:
        snippets = fetch_search_snippets(clean_topic)
        if snippets:
            user_prompt += f"\n\nCurrent context:\n" + "\n".join(f"- {s}" for s in snippets)

    for attempt in range(3):
        try:
            response = openai.ChatCompletion.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.3,
                max_tokens=2000
            )
            content = response['choices'][0]['message']['content'].strip()
            json_start = content.find('{')
            json_end = content.rfind('}') + 1
            if json_start == -1 or json_end == 0:
                continue

            json_str = content[json_start:json_end]
            data = json.loads(json_str)

            slides = []
            for slide in data["slides"]:
                title = slide.get("title", "Untitled").strip()
                content = [c.strip() for c in slide.get("content", []) if c.strip()]
                if title and content:
                    slides.append((title, content))

            if slides:
                return slides[:MAX_SLIDES]
        except (json.JSONDecodeError, KeyError) as e:
            st.warning(f"Retrying generation... (attempt {attempt+1}/3)")
            time.sleep(1)

    raise RuntimeError("Failed to generate valid presentation structure after 3 attempts")

def create_presentation_with_template(slides: List[Tuple[str, List[str]]], template_name: str, filename: str) -> str:
    prs = Presentation(os.path.join(TEMPLATE_DIR, template_name))

    content_layout = None
    for layout in prs.slide_layouts:
        has_title = any(ph.placeholder_format.type == 1 for ph in layout.placeholders)  # TITLE
        has_body = any(ph.placeholder_format.type == 2 for ph in layout.placeholders)   # BODY
        if has_title and has_body:
            content_layout = layout
            break

    if not content_layout:
        raise RuntimeError("No suitable slide layout with title and content placeholders found in the template.")

    if len(prs.slides) == 1:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    for title_text, points in slides:
        slide = prs.slides.add_slide(content_layout)

        if slide.shapes.title:
            slide.shapes.title.text = title_text

        added = False
        for shape in slide.placeholders:
            if shape.is_placeholder and shape.has_text_frame:
                if shape.placeholder_format.type == 2:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    for point in points:
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0
                        p.font.size = Pt(18)
                    added = True
                    break

        if not added:
            for shape in slide.placeholders:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    for point in points:
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0
                        p.font.size = Pt(18)
                    break

    prs.save(filename)
    return filename

st.set_page_config(
    page_title="AI Slide Generator Pro",
    page_icon="📊",
    layout="centered",
    initial_sidebar_state="expanded"
)

with st.sidebar:
    st.header("Configuration")
    selected_model = st.selectbox("AI Model", list(MODELS.keys()), index=0)
    use_live_data = st.checkbox("Include real-time data", True)

    st.markdown("### Choose Slide Template")
    selected_template_label = None
    cols = st.columns(2)
    for i, label in enumerate(template_labels):
        with cols[i % 2]:
            thumb_path = os.path.join(THUMBNAIL_DIR, template_map[label].replace(".pptx", ".png"))
            if os.path.exists(thumb_path):
                st.image(thumb_path, caption=label, use_container_width=True)
            else:
                st.warning(f"No thumbnail for {label}")
            if st.button(f"Use: {label}", key=f"template_{i}"):
                selected_template_label = label

    if not selected_template_label:
        selected_template_label = template_labels[0]
    selected_template = template_map[selected_template_label]

    st.info(f"Using: {MODELS[selected_model]}")
    st.markdown("---")
    st.markdown("**Tips:**")
    st.markdown("- Specify slide count (e.g., '5 slides about AI')")
    st.markdown("- Add 'current data' for latest information")
    st.markdown("- Use quotes for exact topics")

st.title("AI-Powered Presentation Generator")
st.caption("Create professional slides in seconds")

col1, col2 = st.columns([3, 2])
with col1:
    user_prompt = st.text_area(
        "### Presentation Topic",
        height=150,
        placeholder="e.g., Create a 10-slide presentation about Quantum Computing Trends in 2025",
        help="Include 'X-slide' in your prompt to specify slide count"
    )

with col2:
    st.markdown("### Example Prompt")
    st.markdown("- 5-slide intro to quantum computing")
    st.markdown("- 8-slide market analysis: EV industry")
    st.markdown("- 12-slide training: Python for beginners")
    st.markdown("- 7-slide report: climate change impacts")

if st.button("✨ Generate Presentation", type="primary", use_container_width=True):
    if not user_prompt.strip():
        st.warning("Please enter a presentation topic")
    else:
        with st.spinner("Generating your slides..."):
            try:
                status_text = st.empty()
                status_text.text("Analyzing request...")
                time.sleep(0.5)

                status_text.text("Gathering information...")
                slides = generate_slide_content(
                    user_prompt,
                    MODELS[selected_model],
                    use_live_data
                )

                status_text.text("Designing presentation...")
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                    file_path = create_presentation_with_template(slides, selected_template, tmp.name)

                status_text.text("Finalizing...")
                time.sleep(0.5)

                st.success("Presentation generated successfully!")
                st.subheader("Slide Preview")
                for i, (title, points) in enumerate(slides[:3]):
                    with st.expander(f"Slide {i+1}: {title}"):
                        st.markdown("\n".join(f"- {p}" for p in points))

                if len(slides) > 3:
                    st.info(f"+ {len(slides)-3} more slides...")

                with open(file_path, "rb") as f:
                    st.download_button(
                        "📥 Download Presentation",
                        f,
                        file_name=f"presentation_{time.strftime('%Y%m%d')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

            except Exception as e:
                st.error(f"❌ Generation failed: {str(e)}")
                if st.toggle("Show technical details", False):
                    st.exception(e)

st.markdown("---")
st.caption("© 2025 AI Slide Generator Pro | Uses OpenRouter API")
